from abc import ABC, abstractmethod
import os
import re
from typing import Any, Dict, List, Optional, Tuple, Union

import pandas as pd
import pdfplumber
from docx import Document
from lxml import etree
from pptx import Presentation

# ==============================================================    
# BaseTextExtractor
# ==============================================================

class BaseTextExtractor(ABC):
    """
    Abstract base class for text extraction strategies.
    
    This class defines the interface for extracting text from various
    document formats like PDF, Word, PowerPoint, and plain text.
    """
    @abstractmethod
    def extract(self, filepath: Union[str, List[str]]) -> List[str]:
        """
        Extract text from document files.
        
        Args:
            filepath: Path or list of paths to document files
            
        Returns:
            A list of extracted text content, one item per document
        """
        pass

# ==============================================================
# Constants
# ==============================================================

FILE_TYPES = ['pdf', 'ppt', 'pptx', 'pptm', 'doc', 'docx', 'docm', 'txt']

# ==============================================================
# Support Functions
# ==============================================================

# PowerPoint extraction functions
def extract_smartart_text(shape: Any) -> str:
    """
    Recursively extract text from SmartArt shapes in PowerPoint.
    
    Args:
        shape: A PowerPoint shape object
        
    Returns:
        Extracted text from the shape and its children
    """
    text = ""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.text:
            text += paragraph + "\n"
    if shape.shape_type == 16:  # SmartArt type
        for sub_shape in shape.shapes:
            text += extract_smartart_text(sub_shape)
    return text

def ppt_table_to_json(table: Any) -> str:
    """
    Convert a PowerPoint table to JSON format.
    
    Args:
        table: A PowerPoint table object
        
    Returns:
        JSON string representation of the table
    """
    # Extract Content
    newtbl: Dict[int, List[str]] = {}
    for i, row in enumerate(table.rows):
        newtbl[i] = []
        for cell in row.cells:
            # PowerPoint cells don't support nested tables
            newtbl[i].append(cell.text.strip())

    # Check consistent table size
    viable = True
    for i in range(1, len(newtbl)):
        if len(newtbl[i-1]) != len(newtbl[i]):
            viable = False
            break
    
    data = pd.DataFrame(newtbl)
    if viable:
        data.set_index(0, inplace=True)
        data = data.T
        data.reset_index(drop=True, inplace=True)

    # Correct Headers
    heads = []
    for i, c in enumerate(data.columns):
        if c in heads:
            cnt = 0
            new = f"{c}_{cnt}"
            while new in heads:
                new = f"{c}_{cnt}"
                cnt += 1
            data.columns.values[i] = new
            heads.append(new)
        else:
            heads.append(c)

    return data.to_json(index=False)

def ppt_extract(ppt_path: str) -> str:
    """
    Extract text from a PowerPoint file.
    
    Args:
        ppt_path: Path to the PowerPoint file
        
    Returns:
        Extracted text content including slide text, tables, and SmartArt
    """
    # Read PowerPoint file
    ppt = Presentation(ppt_path)

    # Extract text from slides
    slides = []
    for i, slide in enumerate(ppt.slides):
        text_data = ""
        text_data += f"### Slide {i}:\n"
        for shape in slide.shapes:
            if shape.has_table:  # Check for table content
                table = shape.table
                text_data += ppt_table_to_json(table) + "\n" 
            elif shape.shape_type == 16:  # SmartArt type
                text_data += extract_smartart_text(shape)
            else:
                if hasattr(shape, "text"):  # Check shapes
                    text_data += shape.text + "\n"
                elif shape.has_text_frame:  # Check for text boxes
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_data += run.text + "\n"
        # Strip excessive newlines
        while '\n\n\n' in text_data:
            text_data = text_data.replace('\n\n\n', '\n\n')
        slides.append(text_data)

    return '\n\n'.join(slides)

# PDF extraction function
def pdf_extract(pdf_path: str) -> str:
    """
    Extract text from a PDF file.
    
    Args:
        pdf_path: Path to the PDF file
        
    Returns:
        Extracted text content from all pages
    """
    with pdfplumber.open(pdf_path) as pdf:
        text_content = []
        for page in pdf.pages:
            text_content.append(page.extract_text())

        # Join all pages' text into a single string
        full_text = '\n'.join(text_content)

    return full_text 

# Word document extraction functions
def doc_table_to_json(table: Any) -> str:
    """
    Convert a Word document table to JSON format, handling nested tables.
    
    Args:
        table: A Word document table object
        
    Returns:
        JSON string representation of the table
    """
    # Extract Content
    newtbl: Dict[int, List[Any]] = {}
    for i, row in enumerate(table.rows):
        newtbl[i] = []
        for cell in row.cells:
            if len(cell.tables) == 0:
                newtbl[i].append(cell.text.strip())
            else:
                if len(cell.tables) == 1:
                    newtbl[i].append(doc_table_to_json(cell.tables[0]))
                else:
                    nest = []
                    for ctable in cell.tables:
                        nest.append(doc_table_to_json(ctable))
                    newtbl[i].append(','.join(nest))
    
    # Check consistent table size
    viable = True
    for i in range(1, len(newtbl)):
        if len(newtbl[i-1]) != len(newtbl[i]):
            viable = False
            break
    
    data = pd.DataFrame(newtbl)
    if viable:
        data.set_index(0, inplace=True)
        data = data.T
        data.reset_index(drop=True, inplace=True)

    # Correct Headers
    heads = []
    for i, c in enumerate(data.columns):
        if c in heads:
            cnt = 0
            new = f"{c}_{cnt}"
            while new in heads:
                new = f"{c}_{cnt}"
                cnt += 1
            data.columns.values[i] = new
            heads.append(new)
        else:
            heads.append(c)

    return data.to_json(index=False)

def xml_table_to_json(table: Any) -> str:
    """
    Convert an XML table from a Word document to JSON format.
    
    Args:
        table: An XML table element
        
    Returns:
        JSON string representation of the table
    """
    table_xpath = ".//w:txbxContent//w:tbl | .//v:textbox//w:tbl"
    namespaces = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'v': 'urn:schemas-microsoft-com:vml'
    }

    # Extract Content
    newtbl: Dict[int, List[Any]] = {}
    for i, row in enumerate(table.xpath('.//w:tr', namespaces=namespaces)):
        newtbl[i] = []
        for cell in row.xpath('.//w:tc', namespaces=namespaces):
            if len(cell.xpath(table_xpath, namespaces=namespaces)) == 0:
                newtbl[i].append(''.join(cell.xpath('.//w:t//text()', namespaces=namespaces)).strip())
            else:
                if len(cell.xpath(table_xpath, namespaces=namespaces)) == 1:
                    newtbl[i].append(xml_table_to_json(cell.xpath(table_xpath, namespaces=namespaces)[0]))
                else:
                    nest = []
                    for ctable in cell.xpath(table_xpath, namespaces=namespaces):
                        nest.append(xml_table_to_json(ctable))
                    newtbl[i].append(','.join(nest))

    # Check consistent table size
    viable = True
    for i in range(1, len(newtbl)):
        if len(newtbl[i-1]) != len(newtbl[i]):
            viable = False
            break
    
    data = pd.DataFrame(newtbl)
    if viable:
        data.set_index(0, inplace=True)
        data = data.T
        data.reset_index(drop=True, inplace=True)

    # Correct Headers
    heads = []
    for i, c in enumerate(data.columns):
        if c in heads:
            cnt = 0
            new = f"{c}_{cnt}"
            while new in heads:
                new = f"{c}_{cnt}"
                cnt += 1
            data.columns.values[i] = new
            heads.append(new)
        else:
            heads.append(c)

    return data.to_json(index=False)

def extract_textboxes_from_docx(doc: Document) -> List[Dict[str, str]]:
    """
    Extract text boxes from a Word document.
    
    Args:
        doc: A Word document object
        
    Returns:
        List of dictionaries containing text or table content from text boxes
    """
    # Load the document
    text_boxes = []
    duptbl = []
    dupltxt = []

    # Extract raw XML of the document
    xml_content = doc.element.xml
    root = etree.fromstring(xml_content)

    # Define the namespace map
    namespaces = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'v': 'urn:schemas-microsoft-com:vml'
    }

    # XPath to find text boxes
    text_box_xpath = './/w:txbxContent'

    # Search for text boxes using XPath
    for text_box in root.xpath(text_box_xpath, namespaces=namespaces):
        # Extract text within each text box
        text_nodes = text_box.xpath('.//w:t', namespaces=namespaces)
        text_content = ' '.join(''.join(node.itertext()) for node in text_nodes).strip()

        # Check if there is a table within the current text box
        table_xpath = ".//w:tbl"
        tables = text_box.xpath(table_xpath, namespaces=namespaces)

        if tables:
            # Extract tables and format them as JSON
            for table in tables:
                extract = xml_table_to_json(table)
                if extract not in duptbl:  # Avoid duplicates
                    text_boxes.append({"table": extract})
                    duptbl.append(extract)
        else:
            if text_content and text_content not in dupltxt:
                text_boxes.append({"text": text_content})
                dupltxt.append(text_content)

    return text_boxes

def process_text_boxes(doc_path: Document) -> List[str]:
    """
    Process text boxes from a Word document.
    
    Args:
        doc_path: A Word document object
        
    Returns:
        List of extracted text content from text boxes
    """
    text_boxes = extract_textboxes_from_docx(doc_path)

    text_content = []
    xmltext = ""

    for box in text_boxes:
        if "text" in box:
            xmltext += box["text"] + '\n '  # Append text content
        elif "table" in box:
            xmltext += box["table"] + '\n '  # Append table content
            
    text_content.append(xmltext.strip())
    
    return text_content

def doc_extract(doc_path: str) -> str:
    """
    Extract text from a Word document.
    
    Args:
        doc_path: Path to the Word document
        
    Returns:
        Extracted text content including paragraphs, tables, and text boxes
    """
    # Load the .docx file
    doc = Document(doc_path)

    # Extract text from the document
    text_content = []
    for paragraph in doc.paragraphs:
        text_content.append(paragraph.text)

    # Text boxes
    xmltext = ''
    text_boxes = process_text_boxes(doc)
    for text in text_boxes:
        xmltext += text + ' '
    text_content.append(xmltext.strip())

    # Iterate through each table in the document
    for table in doc.tables:
        text_content.append(doc_table_to_json(table))

    # Join all paragraphs into a single string
    full_text = '\n'.join(text_content)

    # Clean up the text
    full_text = re.sub(r'  +', ' ', full_text)
    full_text = full_text.replace(' : ', ': ')
    full_text = full_text.replace(' , ', ', ')
    full_text = full_text.replace("'", "'")
    full_text = full_text.replace(" 's ", "'s ")
    full_text = full_text.replace('"', '"')
    full_text = full_text.replace(" $ ", " $")
    full_text = full_text.replace(" K ", "K ")
    full_text = full_text.replace(" M ", "M ")

    return full_text

# Text document extraction function
def txt_extract(txt_path: str, encodings: Optional[List[str]] = None) -> str:
    """
    Extract text from a text file by trying multiple encodings until one works.
    
    Args:
        txt_path: Path to the text file
        encodings: List of encodings to try. Defaults to common encodings.
    
    Returns:
        The content of the text file
        
    Raises:
        UnicodeDecodeError: If none of the encodings work
    """
    # Default list of encodings to try, in order of likelihood
    if encodings is None:
        encodings = [
            'utf-8', 
            'utf-8-sig',  # UTF-8 with BOM
            'utf-16',
            'utf-16-le',
            'utf-16-be',
            'latin-1',    # Also known as ISO-8859-1, very permissive
            'cp1252',     # Windows default encoding in Western countries
            'ascii'
        ]
    
    # Try each encoding until one works
    for encoding in encodings:
        try:
            with open(txt_path, 'r', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            continue
    
    # If we get here, none of the encodings worked
    raise UnicodeDecodeError(
        f"Failed to decode {txt_path} with any of these encodings: {', '.join(encodings)}"
    )


# ==============================================================
# Class Extractor
# ==============================================================

class TextExtractor(BaseTextExtractor):
    """
    A utility class for extracting text content from various document formats.
    
    Supports PDF, PowerPoint, Word, and plain text files. The class automatically
    detects the file type based on the file extension and uses the appropriate
    extraction method.
    """
    
    def __init__(self) -> None:
        """Initialize the TextExtractor."""
        pass

    def _get_file_type(self, filepath: str) -> str:
        """
        Determine the file type from the file extension.
        
        Args:
            filepath: Path to the file
            
        Returns:
            The lowercase file extension (e.g., 'pdf', 'docx')
        """
        filepath = filepath.lower()
        doc_type = filepath.split('.')[-1]
        return doc_type.lower()
    
    def _extract_powerpoint(self, filepath: str) -> str:
        """
        Extract text from PowerPoint files (.ppt, .pptx, .pptm).
        
        Args:
            filepath: Path to the PowerPoint file
            
        Returns:
            Extracted text content
        """
        print(f"Extracting text from PowerPoint file: {filepath}...")
        return ppt_extract(filepath)

    def _extract_pdf(self, filepath: str) -> str:
        """
        Extract text from PDF files.
        
        Args:
            filepath: Path to the PDF file
            
        Returns:
            Extracted text content
        """
        print(f"Extracting text from PDF file: {filepath}...")
        return pdf_extract(filepath)

    def _extract_word(self, filepath: str) -> str:
        """
        Extract text from Word documents (.doc, .docx, .docm).
        
        Args:
            filepath: Path to the Word document
            
        Returns:
            Extracted text content including paragraphs, tables, and text boxes
        """
        print(f"Extracting text from Word document: {filepath}...")
        return doc_extract(filepath)
    
    def _extract_txt(self, filepath: str) -> str:
        """
        Extract text from plain text files.
        
        Tries multiple encodings to handle different text file formats.
        
        Args:
            filepath: Path to the text file
            
        Returns:
            Extracted text content
            
        Raises:
            UnicodeDecodeError: If the file cannot be decoded with any supported encoding
        """
        print(f"Extracting text from text file: {filepath}...")
        return txt_extract(filepath)
    
    def extract(self, filepath: Union[str, List[str]]) -> List[str]:
        """
        Extract text from document files.
        
        Automatically detects the file type and uses the appropriate extraction method.
        
        Args:
            filepath: Path or list of paths to document files
            
        Returns:
            A list of extracted text content, one item per document
            
        Raises:
            ValueError: If the file type is not supported
        """
        if isinstance(filepath, str):
            filepath = [filepath]

        docs = []
        for file in filepath:
            doc_type = self._get_file_type(file)

            if doc_type not in FILE_TYPES:
                print(f"Filepath: {file} is an unsupported file type: {doc_type}")
                continue

            if 'pdf' in doc_type:
                docs.append(self._extract_pdf(file))
            elif 'ppt' in doc_type:
                docs.append(self._extract_powerpoint(file))
            elif 'doc' in doc_type:
                docs.append(self._extract_word(file))
            elif 'txt' in doc_type:
                docs.append(self._extract_txt(file))
        
        return docs

# ==============================================================
# Other Functions
# ==============================================================

def load_file_paths_from_folder(folder_path: str) -> List[str]:
    """
    Load all file paths from a folder.

    Args:
        folder_path: Path to the folder

    Returns:
        List of file paths
    """
    file_paths = []
    for filename in os.listdir(folder_path):
        file_paths.append(os.path.join(folder_path, filename))
    return file_paths
